from flask import Flask, request, jsonify
from pptx import Presentation
import tempfile
import os

print(">>> PPTX SERVICE FILE LOADED <<<")

app = Flask(__name__)

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"}), 200


def extract_pptx_text(file_path):
    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)

        slides.append({
            "slide": i + 1,
            "text": "\n".join(texts)
        })

    return slides


@app.route("/extract-pptx", methods=["POST"])
def extract_pptx():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        file.save(tmp.name)
        temp_path = tmp.name

    try:
        slides = extract_pptx_text(temp_path)
        return jsonify({
            "success": True,
            "slides": slides
        })

    except Exception as e:
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500

    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


if __name__ == "__main__":
    print(">>> STARTING FLASK SERVER ON PORT 7000 <<<")
    app.run(host="0.0.0.0", port=7000, debug=True)
